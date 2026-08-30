"""Enhanced File Operations - Advanced file system operations for DASH AI OS."""

from __future__ import annotations

import asyncio
import logging
import os
import shutil
import json
import csv
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from datetime import datetime
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class BatchOperationResult:
    """Result of a batch file operation."""
    total: int
    successful: int
    failed: int
    errors: List[str]
    
    def to_dict(self) -> dict:
        return {
            "total": self.total,
            "successful": self.successful,
            "failed": self.failed,
            "errors": self.errors,
        }


class FileOperations:
    """Advanced file system operations.
    
    Features:
    - Batch operations (copy, move, delete)
    - File metadata extraction
    - PDF reading
    - Word document reading
    - Excel reading
    - PowerPoint reading
    - CSV reading
    - JSON reading
    - Code file analysis
    - Folder summaries
    """
    
    def __init__(self):
        self._supported_formats = {
            'pdf': self._read_pdf,
            'docx': self._read_word,
            'doc': self._read_word,
            'xlsx': self._read_excel,
            'xls': self._read_excel,
            'pptx': self._read_powerpoint,
            'ppt': self._read_powerpoint,
            'csv': self._read_csv,
            'json': self._read_json,
        }
    
    async def batch_copy(
        self,
        sources: List[str],
        destination: str,
        overwrite: bool = False,
    ) -> BatchOperationResult:
        """Copy multiple files to a destination.
        
        Args:
            sources: List of source file paths
            destination: Destination directory
            overwrite: Overwrite existing files
            
        Returns:
            BatchOperationResult
        """
        result = BatchOperationResult(
            total=len(sources),
            successful=0,
            failed=0,
            errors=[],
        )
        
        dest_path = Path(destination)
        dest_path.mkdir(parents=True, exist_ok=True)
        
        for source in sources:
            try:
                src_path = Path(source)
                if not src_path.exists():
                    result.failed += 1
                    result.errors.append(f"Source not found: {source}")
                    continue
                
                dest_file = dest_path / src_path.name
                if dest_file.exists() and not overwrite:
                    result.failed += 1
                    result.errors.append(f"Destination exists: {dest_file}")
                    continue
                
                shutil.copy2(src_path, dest_file)
                result.successful += 1
                
            except Exception as e:
                result.failed += 1
                result.errors.append(f"{source}: {str(e)}")
        
        return result
    
    async def batch_move(
        self,
        sources: List[str],
        destination: str,
        overwrite: bool = False,
    ) -> BatchOperationResult:
        """Move multiple files to a destination.
        
        Args:
            sources: List of source file paths
            destination: Destination directory
            overwrite: Overwrite existing files
            
        Returns:
            BatchOperationResult
        """
        result = BatchOperationResult(
            total=len(sources),
            successful=0,
            failed=0,
            errors=[],
        )
        
        dest_path = Path(destination)
        dest_path.mkdir(parents=True, exist_ok=True)
        
        for source in sources:
            try:
                src_path = Path(source)
                if not src_path.exists():
                    result.failed += 1
                    result.errors.append(f"Source not found: {source}")
                    continue
                
                dest_file = dest_path / src_path.name
                if dest_file.exists() and not overwrite:
                    result.failed += 1
                    result.errors.append(f"Destination exists: {dest_file}")
                    continue
                
                shutil.move(str(src_path), str(dest_file))
                result.successful += 1
                
            except Exception as e:
                result.failed += 1
                result.errors.append(f"{source}: {str(e)}")
        
        return result
    
    async def batch_delete(
        self,
        sources: List[str],
    ) -> BatchOperationResult:
        """Delete multiple files.
        
        Args:
            sources: List of file paths to delete
            
        Returns:
            BatchOperationResult
        """
        result = BatchOperationResult(
            total=len(sources),
            successful=0,
            failed=0,
            errors=[],
        )
        
        for source in sources:
            try:
                path = Path(source)
                if not path.exists():
                    result.failed += 1
                    result.errors.append(f"File not found: {source}")
                    continue
                
                if path.is_dir():
                    shutil.rmtree(path)
                else:
                    path.unlink()
                
                result.successful += 1
                
            except Exception as e:
                result.failed += 1
                result.errors.append(f"{source}: {str(e)}")
        
        return result
    
    async def get_file_metadata(self, path: str) -> Dict[str, Any]:
        """Get detailed metadata for a file.
        
        Args:
            path: File path
            
        Returns:
            Metadata dictionary
        """
        try:
            path_obj = Path(path)
            stat = path_obj.stat()
            
            metadata = {
                "name": path_obj.name,
                "path": str(path_obj.absolute()),
                "size": stat.st_size,
                "size_formatted": self._format_size(stat.st_size),
                "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
                "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                "accessed": datetime.fromtimestamp(stat.st_atime).isoformat(),
                "is_directory": path_obj.is_dir(),
                "is_file": path_obj.is_file(),
                "extension": path_obj.suffix,
            }
            
            # Add file-type specific metadata
            if path_obj.is_file():
                ext = path_obj.suffix.lower().lstrip('.')
                if ext in self._supported_formats:
                    specific_metadata = await self._get_specific_metadata(path, ext)
                    metadata.update(specific_metadata)
            
            return metadata
            
        except Exception as e:
            logger.error("Get metadata failed: %s", e)
            return {"error": str(e)}
    
    async def read_file_content(self, path: str) -> Optional[str]:
        """Read file content based on its type.
        
        Args:
            path: File path
            
        Returns:
            File content or None
        """
        try:
            path_obj = Path(path)
            ext = path_obj.suffix.lower().lstrip('.')
            
            if ext in self._supported_formats:
                reader = self._supported_formats[ext]
                return await reader(path)
            else:
                # Try as text
                return path_obj.read_text(encoding='utf-8', errors='ignore')
                
        except Exception as e:
            logger.error("Read file content failed: %s", e)
            return None
    
    async def analyze_code_file(self, path: str) -> Dict[str, Any]:
        """Analyze a code file.
        
        Args:
            path: File path
            
        Returns:
            Analysis results
        """
        try:
            path_obj = Path(path)
            content = path_obj.read_text(encoding='utf-8', errors='ignore')
            
            lines = content.split('\n')
            
            analysis = {
                "language": self._detect_language(path_obj.suffix),
                "total_lines": len(lines),
                "code_lines": len([l for l in lines if l.strip() and not l.strip().startswith('#')]),
                "comment_lines": len([l for l in lines if l.strip().startswith('#')]),
                "blank_lines": len([l for l in lines if not l.strip()]),
                "functions": self._extract_functions(content, path_obj.suffix),
                "classes": self._extract_classes(content, path_obj.suffix),
                "imports": self._extract_imports(content, path_obj.suffix),
            }
            
            return analysis
            
        except Exception as e:
            logger.error("Code analysis failed: %s", e)
            return {"error": str(e)}
    
    async def generate_folder_summary(self, path: str) -> Dict[str, Any]:
        """Generate a summary of a folder.
        
        Args:
            path: Folder path
            
        Returns:
            Folder summary
        """
        try:
            path_obj = Path(path)
            if not path_obj.is_dir():
                return {"error": "Not a directory"}
            
            files = []
            folders = []
            total_size = 0
            
            for item in path_obj.iterdir():
                try:
                    stat = item.stat()
                    total_size += stat.st_size
                    
                    if item.is_dir():
                        folders.append({
                            "name": item.name,
                            "items": len(list(item.iterdir())),
                        })
                    else:
                        files.append({
                            "name": item.name,
                            "size": stat.st_size,
                            "type": item.suffix,
                        })
                except (PermissionError, OSError):
                    continue
            
            return {
                "name": path_obj.name,
                "path": str(path_obj.absolute()),
                "total_items": len(files) + len(folders),
                "file_count": len(files),
                "folder_count": len(folders),
                "total_size": total_size,
                "total_size_formatted": self._format_size(total_size),
                "files": files[:20],  # Limit to first 20
                "folders": folders[:20],
            }
            
        except Exception as e:
            logger.error("Folder summary failed: %s", e)
            return {"error": str(e)}
    
    # ── File Type Readers ─────────────────────────────────────
    
    async def _read_pdf(self, path: str) -> Optional[str]:
        """Read PDF file content."""
        try:
            import PyPDF2
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except ImportError:
            logger.warning("PyPDF2 not installed")
            return None
        except Exception as e:
            logger.error("PDF read failed: %s", e)
            return None
    
    async def _read_word(self, path: str) -> Optional[str]:
        """Read Word document content."""
        try:
            import docx
            doc = docx.Document(path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text
        except ImportError:
            logger.warning("python-docx not installed")
            return None
        except Exception as e:
            logger.error("Word read failed: %s", e)
            return None
    
    async def _read_excel(self, path: str) -> Optional[str]:
        """Read Excel file content."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path)
            text = ""
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                for row in sheet.iter_rows(values_only=True):
                    text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                text += "\n"
            return text
        except ImportError:
            logger.warning("openpyxl not installed")
            return None
        except Exception as e:
            logger.error("Excel read failed: %s", e)
            return None
    
    async def _read_powerpoint(self, path: str) -> Optional[str]:
        """Read PowerPoint file content."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text
        except ImportError:
            logger.warning("python-pptx not installed")
            return None
        except Exception as e:
            logger.error("PowerPoint read failed: %s", e)
            return None
    
    async def _read_csv(self, path: str) -> Optional[str]:
        """Read CSV file content."""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
                text = "\n".join([",".join(row) for row in rows])
                return text
        except Exception as e:
            logger.error("CSV read failed: %s", e)
            return None
    
    async def _read_json(self, path: str) -> Optional[str]:
        """Read JSON file content."""
        try:
            with open(path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, indent=2)
        except Exception as e:
            logger.error("JSON read failed: %s", e)
            return None
    
    # ── Helper Methods ───────────────────────────────────────
    
    def _format_size(self, size: int) -> str:
        """Format file size."""
        for unit in ['B', 'KB', 'MB', 'GB', 'TB']:
            if size < 1024:
                return f"{size:.1f} {unit}"
            size /= 1024
        return f"{size:.1f} PB"
    
    async def _get_specific_metadata(self, path: str, ext: str) -> Dict[str, Any]:
        """Get file-type specific metadata."""
        metadata = {}
        
        if ext == 'pdf':
            metadata.update(await self._get_pdf_metadata(path))
        elif ext in ['docx', 'doc']:
            metadata.update(await self._get_word_metadata(path))
        elif ext in ['xlsx', 'xls']:
            metadata.update(await self._get_excel_metadata(path))
        
        return metadata
    
    async def _get_pdf_metadata(self, path: str) -> Dict[str, Any]:
        """Get PDF-specific metadata."""
        try:
            import PyPDF2
            with open(path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                return {
                    "pages": len(reader.pages),
                    "title": reader.metadata.get('/Title', '') if reader.metadata else '',
                    "author": reader.metadata.get('/Author', '') if reader.metadata else '',
                }
        except Exception:
            return {}
    
    async def _get_word_metadata(self, path: str) -> Dict[str, Any]:
        """Get Word-specific metadata."""
        try:
            import docx
            doc = docx.Document(path)
            return {
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
            }
        except Exception:
            return {}
    
    async def _get_excel_metadata(self, path: str) -> Dict[str, Any]:
        """Get Excel-specific metadata."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path)
            sheet_info = {}
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_info[sheet_name] = {
                    "rows": sheet.max_row,
                    "columns": sheet.max_column,
                }
            return {
                "sheets": len(wb.sheetnames),
                "sheet_info": sheet_info,
            }
        except Exception:
            return {}
    
    def _detect_language(self, extension: str) -> str:
        """Detect programming language from extension."""
        lang_map = {
            '.py': 'Python',
            '.js': 'JavaScript',
            '.ts': 'TypeScript',
            '.html': 'HTML',
            '.css': 'CSS',
            '.java': 'Java',
            '.cpp': 'C++',
            '.c': 'C',
            '.h': 'C/C++ Header',
            '.go': 'Go',
            '.rs': 'Rust',
            '.swift': 'Swift',
            '.rb': 'Ruby',
            '.php': 'PHP',
            '.sh': 'Shell',
            '.sql': 'SQL',
        }
        return lang_map.get(extension.lower(), 'Unknown')
    
    def _extract_functions(self, content: str, extension: str) -> List[str]:
        """Extract function names from code."""
        # Simplified implementation
        if extension in ['.py']:
            import re
            return re.findall(r'def\s+(\w+)', content)
        elif extension in ['.js', '.ts']:
            import re
            return re.findall(r'function\s+(\w+)|(\w+)\s*=\s*\([^)]*\)\s*=>', content)
        return []
    
    def _extract_classes(self, content: str, extension: str) -> List[str]:
        """Extract class names from code."""
        if extension in ['.py']:
            import re
            return re.findall(r'class\s+(\w+)', content)
        elif extension in ['.js', '.ts']:
            import re
            return re.findall(r'class\s+(\w+)', content)
        return []
    
    def _extract_imports(self, content: str, extension: str) -> List[str]:
        """Extract imports from code."""
        if extension in ['.py']:
            import re
            return re.findall(r'import\s+(\w+)|from\s+(\w+)', content)
        elif extension in ['.js', '.ts']:
            import re
            return re.findall(r'import.*from\s+[\'"]([^\'"]+)[\'"]', content)
        return []


_file_operations: Optional[FileOperations] = None


def get_file_operations() -> FileOperations:
    global _file_operations
    if _file_operations is None:
        _file_operations = FileOperations()
    return _file_operations
